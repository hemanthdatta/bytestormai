import os
import io
import re
import traceback
import google.generativeai as genai
from PIL import Image
import pandas as pd
from collections import Counter

# For DOCX processing
from docx import Document
from docx.document import Document as DocumentClass
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table, _Cell
from docx.text.paragraph import Paragraph

# For PPTX processing
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pdf2md import summarize_markdown_groups

def create_safe_filename(text):
    """Create a safe filename from text by removing invalid characters."""
    return re.sub(r'[\\/*?:"<>|]', "", text.replace(' ', '_'))[:50]

def get_font_size_stats(presentation):
    """
    Scan the entire PPTX to collect and count all font sizes.
    """
    font_sizes = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run, "font") and hasattr(run.font, "size") and run.font.size:
                            font_sizes.append(round(run.font.size / 100))  # Convert EMUs to points (approximate)
    return Counter(font_sizes)

def get_heading_thresholds(font_counter):
    """
    Determine thresholds for section (##) and subsection (###) headings.
    Most frequent size = body text; largest above that = section; next = subsection.
    """
    if not font_counter:
        return 0, 0
    sizes = sorted(font_counter.keys(), reverse=True)
    body_size = font_counter.most_common(1)[0][0]
    section_size = sizes[0] if sizes[0] > body_size else body_size + 1
    subsection_size = (sizes[1] if len(sizes) > 1 and sizes[1] > body_size else body_size + 0.5)
    return section_size, subsection_size

def generate_md(pptx_path, api_key, output_file=None, max_workers=4):
    """
    Extract content from a PPTX file into a markdown file with images and text.
    
    Args:
        pptx_path (str): Path to the PPTX file
        api_key (str): Google API key for Gemini
        output_file (str): Optional output file name (without .md extension)
        max_workers (int): Maximum number of parallel workers (not used for PPTX, kept for API compatibility)
    """
    try:
        # Configure the Gemini API
        genai.configure(api_key=api_key)
        
        # Check if file exists
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PPTX file not found: {pptx_path}")
            
        # Setup output directory
        pptx_name = os.path.splitext(os.path.basename(pptx_path))[0]
        if output_file is None:
            output_file = pptx_name
            output_dir = os.path.join(os.path.dirname(pptx_path), f"{pptx_name}_markdown")
        else:
            output_dir = os.path.dirname(os.path.abspath(output_file))
            if not output_dir:  # If output_file has no directory component
                output_dir = os.path.join(os.path.dirname(pptx_path), f"{output_file}_markdown")
        
        # Create output directories
        images_dir = os.path.join(output_dir, "images")
        slide_images_dir = os.path.join(output_dir, "slide_images")
        os.makedirs(images_dir, exist_ok=True)
        os.makedirs(slide_images_dir, exist_ok=True)
        
        # Initialize markdown content
        markdown_content = f"# {pptx_name}\n\n"
        
        # Initialize Gemini model - with error handling
        try:
            model = genai.GenerativeModel('gemini-2.0-flash')
        except Exception as e:
            print(f"Error initializing Gemini model: {str(e)}")
            print("Trying with model name 'gemini-pro-vision'...")
            try:
                model = genai.GenerativeModel('gemini-pro-vision')
            except Exception as e2:
                print(f"Error with alternate model: {str(e2)}")
                raise Exception("Failed to initialize Gemini model. Check your API key and network connection.")
        
        # Open the PPTX with error handling
        try:
            presentation = Presentation(pptx_path)
        except Exception as e:
            raise Exception(f"Failed to open PPTX: {str(e)}")
        
        # Get font statistics and heading thresholds
        font_counter = get_font_size_stats(presentation)
        section_size, subsection_size = get_heading_thresholds(font_counter)
        print(f"PPTX analysis: section_size={section_size}, subsection_size={subsection_size}")
        
        # Process each slide
        for slide_num, slide in enumerate(presentation.slides, start=1):
            try:
                # Add page/slide header with horizontal rule for separation
                markdown_content += f"\n\n---\n\n## Page {slide_num}\n\n"
                
                # Add slide reference
                markdown_content += f"*Slide {slide_num} of {len(presentation.slides)}*\n\n"
                
                # Extract slide title if available
                if slide.shapes.title and hasattr(slide.shapes.title, "text_frame") and slide.shapes.title.text_frame:
                    title_text = slide.shapes.title.text_frame.text
                    if title_text.strip():
                        markdown_content += f"### {title_text.strip()}\n\n"
                
                # Process text content
                text_content = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        # Skip if this is the title we already processed
                        if shape == slide.shapes.title:
                            continue
                            
                        # Process paragraphs to check for headings based on font size
                        for paragraph in shape.text_frame.paragraphs:
                            # Skip empty paragraphs
                            if not paragraph.text.strip():
                                continue
                                
                            # Get maximum font size in this paragraph
                            max_size = 0
                            for run in paragraph.runs:
                                if hasattr(run, "font") and hasattr(run.font, "size") and run.font.size:
                                    size = run.font.size / 100  # Convert EMUs to points (approximate)
                                    max_size = max(max_size, size)
                            
                            # Format text based on font size
                            if max_size >= section_size:
                                text_content.append((2, paragraph.text.strip()))  # Level 2 heading
                            elif max_size >= subsection_size:
                                text_content.append((3, paragraph.text.strip()))  # Level 3 heading
                            else:
                                text_content.append((0, paragraph.text.strip()))  # Regular text
                
                # Add text content if any
                if text_content:
                    for level, text in text_content:
                        if level == 2:
                            markdown_content += f"## {text}\n\n"
                        elif level == 3:
                            markdown_content += f"### {text}\n\n"
                        else:
                            markdown_content += f"{text}\n\n"
                
                # Process images
                image_count = 0
                
                for shape_idx, shape in enumerate(slide.shapes):
                    # Check if shape is a picture
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image_count += 1
                            
                            # Save the image
                            img_filename = f"slide_{slide_num}_image_{image_count}.png"
                            img_path = os.path.join(images_dir, img_filename)
                            
                            # Extract image
                            image = shape.image
                            image_bytes = image.blob
                            
                            with open(img_path, "wb") as img_file:
                                img_file.write(image_bytes)
                            
                            # Add image to markdown
                            markdown_content += f"### Image {image_count}\n\n"
                            markdown_content += f"![Image {image_count} on Slide {slide_num}](images/{img_filename})\n\n"
                            
                            # Use Gemini for OCR on the image
                            try:
                                pil_image = Image.open(io.BytesIO(image_bytes))
                                response = model.generate_content(["Extract all text from this image", pil_image])
                                ocr_text = response.text
                                if ocr_text and ocr_text.strip():
                                    markdown_content += f"**OCR Text from Image {image_count}:**\n\n{ocr_text}\n\n"
                            except Exception as e:
                                markdown_content += f"**OCR Failed for Image {image_count}:** {str(e)}\n\n"
                                
                        except Exception as e:
                            markdown_content += f"**Error processing Image {image_count} on Slide {slide_num}:** {str(e)}\n\n"
                
                # Try to save whole slide as image for reference
                try:
                    # We can't directly save the slide as an image using python-pptx
                    # For now, we'll note this as a limitation
                    markdown_content += "Note: Full slide image not available with python-pptx library\n\n"
                except Exception as e:
                    print(f"Error saving slide image: {str(e)}")
                
                print(f"Processed slide {slide_num} of {len(presentation.slides)}")
                
            except Exception as e:
                markdown_content += f"**Error processing Slide {slide_num}:** {str(e)}\n\n---\n\n"
                print(f"Error processing slide {slide_num}: {str(e)}")
        
        # Save the final markdown file
        markdown_file = f"{output_file}.md"
        with open(markdown_file, "w", encoding="utf-8") as f:
            f.write(markdown_content)
        
        print(f"Markdown file saved to: {markdown_file}")
        
    except Exception as e:
        print(f"Error occurred: {str(e)}")
        print(traceback.format_exc())
